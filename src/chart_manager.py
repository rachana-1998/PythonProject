from pptx.chart import chart
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import io
import logging
from typing import Literal, Union, List, Dict, Any
logger = logging.getLogger('mcp_chart_manager')

# Theme definitions (copied from presentation_manager.py for consistency)
THEMES = {
    "modern_blue": {
        "background": (0x00, 0x5A, 0xC1),
        "text": (0xFF, 0xFF, 0xFF),
        "accent": (0xE0, 0xF7, 0xFA),
        "font": "Montserrat"
    },
    "elegant_green": {
        "background": (0x2E, 0x7D, 0x32),
        "text": (0xFF, 0xFF, 0xFF),
        "accent": (0xC8, 0xE6, 0xC9),
        "font": "Lato"
    }
}

class ChartManager:
    def __init__(self):
        self.name = "Chart Manager"

    def determine_chart_type(self, data: Dict[str, Any]) -> tuple[XL_CHART_TYPE, str]:
        """
        Analyze the data structure and determine the most appropriate chart type.
        Returns tuple of (PowerPoint chart type enum, chart_format)
        """
        try:
            series_count = len(data.get("series", []))
            categories = data.get("categories", [])

            is_xy_data = False
            for series in data.get("series", []):
                values = series.get("values", [])
                if values and isinstance(values[0], (list, tuple)) and len(values[0]) == 2:
                    is_xy_data = True
                    break

            if is_xy_data:
                logger.debug("Selected XY_SCATTER chart type")
                return XL_CHART_TYPE.XY_SCATTER, "xy"

            # Pie chart for single series with percentage-like data
            if series_count == 1 and categories:
                values = data["series"][0].get("values", [])
                if len(values) <= 8:
                    try:
                        total = sum(float(v) for v in values)
                        if 95 <= total <= 105:
                            logger.debug("Selected PIE chart type")
                            return XL_CHART_TYPE.PIE, "category"
                    except (TypeError, ValueError):
                        pass

                # Line chart for time-based categories
                if categories and any(
                    isinstance(cat, (str, int)) and
                    any(term in str(cat).lower() for term in
                        ["date", "time", "year", "month", "quarter", "q1", "q2", "q3", "q4"])
                    for cat in categories
                ):
                    logger.debug("Selected LINE chart type")
                    return XL_CHART_TYPE.LINE, "category"

            # Bar chart for multiple series with categories
            if series_count > 1 and categories:
                logger.debug("Selected BAR_CLUSTERED chart type")
                return XL_CHART_TYPE.BAR_CLUSTERED, "category"

            # Default to column chart
            logger.debug("Selected COLUMN_CLUSTERED chart type (default)")
            return XL_CHART_TYPE.COLUMN_CLUSTERED, "category"

        except Exception as e:
            logger.error(f"Error determining chart type: {str(e)}")
            raise ValueError(f"Invalid chart data: {str(e)}")

    def add_chart_to_slide(self, slide, chart_type: XL_CHART_TYPE, data: Dict[str, Any],
                           chart_format: str = "category", theme: str = "modern_blue") -> chart:
        """
        Add a chart to the slide with the specified data and theme styling.
        Optionally uses matplotlib for enhanced rendering if chart_type is not supported.
        """
        theme_data = THEMES.get(theme, THEMES["modern_blue"])
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(5)

        # Use matplotlib for advanced charts if desired
        if chart_type in [XL_CHART_TYPE.BAR_CLUSTERED, XL_CHART_TYPE.LINE, XL_CHART_TYPE.PIE]:
            try:
                plt.figure(figsize=(8, 5))
                categories = data.get("categories", [])
                series = data.get("series", [])

                if chart_type == XL_CHART_TYPE.PIE:
                    values = series[0]["values"] if series else []
                    plt.pie(values, labels=categories, colors=[
                        f"#{theme_data['accent'][0]:02x}{theme_data['accent'][1]:02x}{theme_data['accent'][2]:02x}"
                    ])
                    plt.title(data.get("title", ""))
                elif chart_type == XL_CHART_TYPE.LINE:
                    for s in series:
                        plt.plot(categories, s["values"], label=s["name"])
                    plt.legend()
                    plt.xlabel(data.get("x_axis", ""))
                    plt.ylabel(data.get("y_axis", ""))
                else:  # BAR_CLUSTERED
                    for s in series:
                        plt.bar(categories, s["values"], label=s["name"],
                                color=f"#{theme_data['accent'][0]:02x}{theme_data['accent'][1]:02x}{theme_data['accent'][2]:02x}")
                    plt.legend()
                    plt.xlabel(data.get("x_axis", ""))
                    plt.ylabel(data.get("y_axis", ""))

                img_stream = io.BytesIO()
                plt.savefig(img_stream, format="png", bbox_inches="tight")
                plt.close()
                slide.shapes.add_picture(img_stream, left, top, width, height)
                logger.info(f"Added matplotlib {chart_type.name} chart to slide")
                return None  # Added as image, not pptx chart
            except Exception as e:
                logger.warning(f"Matplotlib rendering failed: {str(e)}. Falling back to python-pptx chart.")

        # Default python-pptx chart rendering
        try:
            if chart_format == "category":
                chart_data = CategoryChartData()
                chart_data.categories = data.get("categories", [])
                for series in data["series"]:
                    chart_data.add_series(series["name"], series["values"])
            elif chart_format == "xy":
                chart_data = XyChartData()
                for series in data["series"]:
                    series_data = chart_data.add_series(series["name"])
                    for x, y in series["values"]:
                        series_data.add_data_point(x, y)

            graphic_frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
            chart_obj = graphic_frame.chart

            # Apply theme styling
            chart_obj.has_legend = True
            if len(data["series"]) > 1:
                chart_obj.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart_obj.legend.font.name = theme_data["font"]
                chart_obj.legend.font.size = Pt(10)
                chart_obj.legend.font.color.rgb = RGBColor(*theme_data["text"])

            if "x_axis" in data:
                chart_obj.category_axis.axis_title.text_frame.text = data["x_axis"]
                p = chart_obj.category_axis.axis_title.text_frame.paragraphs[0]
                p.font.name = theme_data["font"]
                p.font.color.rgb = RGBColor(*theme_data["text"])
                p.font.size = Pt(12)

            if "y_axis" in data:
                chart_obj.value_axis.axis_title.text_frame.text = data["y_axis"]
                p = chart_obj.value_axis.axis_title.text_frame.paragraphs[0]
                p.font.name = theme_data["font"]
                p.font.color.rgb = RGBColor(*theme_data["text"])
                p.font.size = Pt(12)

            # Apply accent color to series
            for series in chart_obj.series:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = RGBColor(*theme_data["accent"])
                if series.data_labels:
                    series.data_labels.font.name = theme_data["font"]
                    series.data_labels.font.color.rgb = RGBColor(*theme_data["text"])

            logger.info(f"Added {chart_type.name} chart to slide with theme '{theme}'")
            return chart_obj

        except Exception as e:
            logger.error(f"Failed to add chart: {str(e)}")
            raise ValueError(f"Failed to add chart: {str(e)}")
# return chart