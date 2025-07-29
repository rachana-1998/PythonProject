import os
import json
import aiohttp
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Inches
from pptx.slide import Slide
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from PIL import Image, UnidentifiedImageError

import logging
from typing import Literal, Union, List, Dict, Any
logger = logging.getLogger('mcp_presentation_manager')


# initialize AI client
async_anthropic = os.getenv("HF_API_Token")
segmind_api_key = os.getenv("SM_API_Key")

ChartTypes = Literal["bar", "line", "pie", "scatter", "area"]


# Layout templates
LAYOUTS = {
    "title": {
        "placeholders": ["title", "subtitle"],
        "style": {"bg_color": "RGB(0, 90, 193)", "font": "Montserrat"}
    },
    "content": {
        "placeholders": ["title", "bullet_points"],
        "style": {"bg_gradient": ["RGB(0, 90, 193)", "RGB(224, 247, 250)"]}
    },
    "chart": {
        "placeholders": ["title", "chart"],
        "style": {"bg_color": "RGB(245, 245, 245)"}
    },
    "image": {
        "placeholders": ["title", "image"],
        "style": {"bg_color": "RGB(224, 224, 224)"}
    }
}

# Theme definitions
THEMES = {
    "modern_blue": {
        "background": (0x00, 0x5A, 0xC1),
        "text": (0xFF, 0xFF, 0xFF),
        "accent": (0xE0, 0xF7, 0xFA),
        "font": "Montserrat"
    },
    "elegant_green": {
        "background": (0x2E, 0x7D, 0x32),
        "text": (0xFF, 0xFF, 0xFF),
        "accent": (0xC8, 0xE6, 0xC9),
        "font": "Lato"
    }
}
# json structure
ClaudeJson ={
  "slides": [
    {
      "type": "title",
      "title": "Project Overview",
      "subtitle": "Q3 2025 Update",
      "theme": "gradient_blue"
    },
    {
      "type": "content",
      "title": "Key Points",
      "bullets": ["Point 1", "Point 2", "Point 3"],
      "image_prompt": "Modern office setting",
      "theme": "minimalist"
    },
    {
      "type": "table",
      "title": "Sales Data",
      "data": [["Region", "Sales"], ["North", 1000], ["South", 1500]],
      "theme": "professional"
    }
  ]
}

class PresentationManager:
    # Slide layout constants
    SLIDE_LAYOUT_TITLE = 0
    SLIDE_LAYOUT_TITLE_AND_CONTENT = 1
    SLIDE_LAYOUT_SECTION_HEADER = 2
    SLIDE_LAYOUT_TWO_CONTENT = 3
    SLIDE_LAYOUT_COMPARISON = 4
    SLIDE_LAYOUT_TITLE_ONLY = 5
    SLIDE_LAYOUT_BLANK = 6
    SLIDE_LAYOUT_CONTENT_WITH_CAPTION = 7
    SLIDE_LAYOUT_PICTURE_WITH_CAPTION = 8




    def __init__(self):
        self.presentations: Dict[str, Any] = {}


    def apply_theme(self, presentation: Presentation, theme) -> dict:
        """Apply a theme to the presentation's slide master."""
        theme = THEMES.get(theme, THEMES["modern_blue"])

        slide_master = presentation.slide_master
        bg = slide_master.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*theme["background"])
        logger.info(f"Applied theme {theme} to presentation")
        return theme

    prompt = """
        Generate a JSON outline for a 5-slide PowerPoint presentation about sustainable architecture. Include:
        - slide_title: string
        - content: list of strings or objects (for tables/charts)
        - layout: string (title, content, chart, image_text, table)
        - styling: object with color_scheme (primary/secondary colors), font, and image_prompt (for Stable Diffusion)
        Ensure the response is valid JSON.
        """

    def get_presentation_structure(prompt, model="claude-3.5-sonnet"):
        """Query Claude for presentation structure in JSON format."""
        client = async_anthropic
        response = client.messages.create(
            model=model,
            max_tokens=1000,
            messages=[{"role": "user", "content": prompt}]
        )
        # Parse JSON from Claude's response
        try:
            return json.loads(response.content[0].text)
        except json.JSONDecodeError:
            raise ValueError("Claude did not return valid JSON")

    async def generate_image(prompt):
        """Generate image using Segmind API"""
        async with aiohttp.ClientSession() as session:
            async with session.post(
                    "https://api.segmind.com/v1/stable-diffusion",
                    headers={"SM_API_Key": segmind_api_key},
                    json={
                        "prompt": prompt,
                        "negative_prompt": "low quality, blurry",
                        "samples": 1,
                        "scheduler": "DPM++ 2M",
                        "num_inference_steps": 25,
                        "guidance_scale": 7.5,
                        "seed": 1234,
                        "img_width": 512,
                        "img_height": 512
                    }
            ) as response:
                if response.status == 200:
                    with open("temp_image.jpg", "wb") as f:
                        f.write(await response.read())
                    return "temp_image.jpg"
                return None
    def create_presentation(self, presentation_name: str,json_input: str=None, theme: str = "modern_blue") -> Presentation:
        """Create a new presentation with the specified theme."""
        prs = Presentation()
        self.apply_theme(prs, theme)
        # Parse JSON input if provided
        if json_input:
            try:
                data = json.loads(json_input)
                if not isinstance(data, dict) or "slides" not in data:
                    raise ValueError("Invalid JSON format. Expected a 'slides' key.")

                for slide_data in data["slides"]:
                    slide_type = slide_data.get("type")
                    slide_theme = slide_data.get("theme", theme)  # Use slide-specific theme or default
                    if slide_type == "title":
                        self.add_title_slide(prs, slide_data, slide_theme)
                    elif slide_type == "content":
                        self.add_content_slide(prs, slide_data, slide_theme)
                    elif slide_type == "table":
                        self.add_table_slide(prs, slide_data, slide_theme)
                    else:
                        logger.warning(f"Unsupported slide type '{slide_type}'")
            except json.JSONDecodeError:
                logger.error("Failed to parse JSON input")
                raise ValueError("Invalid JSON input")
            except Exception as e:
                logger.error(f"Error processing JSON: {str(e)}")
                raise

        self.presentations[presentation_name] = prs
        logger.info(f"Created presentation {presentation_name} with theme {theme}")
        return prs

    def add_formatted_bullets(self, text_frame, text_block, theme: dict):
        """
        Process a text block and add paragraphs with proper bullet indentation
        using ASCII code detection:
        - ASCII 10 (LF) or ASCII 13 (CR) or combination for new lines (main bullets)
        - ASCII 9 (HT) for tab indentation (sub-bullets)

        Args:
            text_frame: The PowerPoint text frame to add text to
            text_block: String of text to process
        """
        # First, normalize all line endings to a single format
        # Replace CR+LF (Windows) with a single marker
        normalized_text = text_block.replace('\r\n', '\n').replace('\r', '\n')
        # Replace any remaining CR (old Mac) with LF
        #normalized_text = normalized_text.replace('\r', '\n')

        # Split the text block into lines using ASCII 10 (LF)
        lines = normalized_text.split('\n')

        # Clear any existing text
        if text_frame.paragraphs:
            p = text_frame.paragraphs[0]
            p.text = ""
        else:
            p = text_frame.add_paragraph()

        # # Process the first line separately (if it exists)
        # if lines and lines[0].strip():
        #     first_line = lines[0]
        #     # Count leading tabs (ASCII 9) to determine indentation level
        #     level = 0
        #     while first_line and ord(first_line[0]) == 9:  # ASCII 9 is HT (tab)
        #         level += 1
        #         first_line = first_line[1:]
        #
        #     p.text = first_line.strip()
        #     p.level = level
        #
        # # Process remaining lines
        for i, line in enumerate(lines):
            if not line.strip():
                continue  # Skip empty lines

            # Count leading tabs (ASCII 9) to determine indentation level
            level = 0
            while line and ord(line[0]) == 9:  # ASCII 9 is HT (tab)
                level += 1
                line = line[1:]

            # Add the paragraph with proper indentation
            p = text_frame.add_paragraph() if i > 0 else p
            p.text = line.strip()
            p.level = level
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(*theme["text"])
            p.font.name = theme["font"]
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.word_wrap = True

    def add_section_header_slide(self, presentation_name: str, header: str, subtitle: str = None, theme: str = "modern_blue") -> Slide:
        """
        Create a section header slide for the given presentation

        Args:
            presentation_name: The presentation to add the slide to
            header: The section header to use
            subtitle: The subtitle of the section header to use
        """
        try:
            prs = self.presentations[presentation_name]
        except KeyError as e:
            logger.error(f"Presentation '{presentation_name}' not found")
            raise ValueError(f"Presentation '{presentation_name}' not found")
        #slide_master = prs.slide_master
        theme_data = THEMES.get(theme, THEMES["modern_blue"])
        # Add a new slide with layout
        slide_layout = prs.slide_layouts[self.SLIDE_LAYOUT_SECTION_HEADER]
        slide = prs.slides.add_slide(slide_layout)

        # # Set the subtitle
        # if subtitle:
        #     subtitle_shape = slide.placeholders[1]
        #     text_frame = subtitle_shape.text_frame
        #     text_frame.text = subtitle
        #
        # # Set the section header
        if header:
            header_shape = slide.shapes.title
            header_shape.text = header
            header_shape.text_frame.paragraphs[0].font.size = Pt(32)
            header_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*theme_data["text"])
            header_shape.text_frame.paragraphs[0].font.name = theme_data["font"]

        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text_frame.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(18)
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*theme_data["text"])
            subtitle_shape.text_frame.paragraphs[0].font.name = theme_data["font"]

        logger.info(f"Added section header slide '{header}' to {presentation_name}")

        return slide

    def add_comparison_slide(self, presentation_name: str, title: str, left_side_title: str, left_side_content: str,
                             right_side_title: str, right_side_content: str, theme: str = "modern_blue") -> Slide:
        """

        Create a comparison slide with theme styling.
        Args:
            presentation_name: The presentation to add the slide to
            title: The title of the slide
            left_side_title: The title of the left hand side content
            left_side_content: The body content for the left hand side
            right_side_title: The title of the right hand side content
            right_side_content: The body content for the right hand side
        """
        try:
            prs = self.presentations[presentation_name]
        except KeyError as e:
            logger.error(f"Presentation '{presentation_name}' not found")
            raise ValueError(f"Presentation '{presentation_name}' not found")
        #slide_master = prs.slide_master

        theme_data = THEMES.get(theme, THEMES["modern_blue"])

        # Add a new slide with layout
        slide_layout = prs.slide_layouts[self.SLIDE_LAYOUT_COMPARISON]
        slide = prs.slides.add_slide(slide_layout)

        # Set the title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*theme_data["text"])
        title_shape.text_frame.paragraphs[0].font.name = theme_data["font"]

        # Build the left hand content
        content_shape = slide.placeholders[1]
        content_shape.text_frame.text = left_side_title
        content_shape.text_frame.paragraphs[0].font.size = Pt(18)
        content_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*theme_data["text"])
        content_shape.text_frame.paragraphs[0].font.name = theme_data["font"]
        content_shape = slide.placeholders[2]
        self._add_formatted_bullets(content_shape.text_frame, left_side_content, theme_data)

        # Build the right hand content
        content_shape = slide.placeholders[3]
        content_shape.text_frame.text = right_side_title
        content_shape.text_frame.paragraphs[0].font.size = Pt(18)
        content_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*theme_data["text"])
        content_shape.text_frame.paragraphs[0].font.name = theme_data["font"]
        content_shape = slide.placeholders[4]
        self._add_formatted_bullets(content_shape.text_frame, right_side_content, theme_data)

        logger.info(f"Added comparison slide '{title}' to {presentation_name}")
        return slide

    def add_picture_with_caption_slide(self, presentation_name: str, title: str, image_path: str, caption_text: str,
                                       theme: str = "modern_blue") -> Slide:
        """
        For the given presentation builds a slide with the picture with caption template.
        Maintains the image's aspect ratio by adjusting the picture object after insertion.
        Args:
            presentation_name: The presentation to add the slide to
            title: The title of the slide
            image_path: The path to the image to insert
            caption_text: The caption content

        """
        try:
            prs = self.presentations[presentation_name]
        except KeyError as e:
            logger.error(f"Presentation '{presentation_name}' not found")
            raise ValueError(f"Presentation '{presentation_name}' not found")

            theme_data = THEMES.get(theme, THEMES["modern_blue"])
            slide_layout = prs.slide_layouts[self.SLIDE_LAYOUT_PICTURE_WITH_CAPTION]
            slide = prs.slides.add_slide(slide_layout)
        except IndexError as e:
            error_message = f"Slide Index does not exist. Error: {str(e)}"
            raise ValueError(error_message)
        # Set the title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*theme_data["text"])
        title_shape.text_frame.paragraphs[0].font.name = theme_data["font"]

        # Get the image placeholder
        try:
            placeholder = slide.placeholders[1]
        except IndexError as e:
            logger.error(f"Placeholder index does not exist. Error: {str(e)}")
            raise ValueError(f"Placeholder index does not exist. Error: {str(e)}")

        # Insert the picture into the placeholder
        if not os.path.exists(image_path):
            raise FileNotFoundError(f"Image not found: {image_path}")
        try:
            img = Image.open(image_path).convert("RGB").resize((800, 600)).enhance(1.2)
            processed_path = os.path.join(os.path.dirname(image_path), "processed_image.png")
            img.save(processed_path)
            picture = placeholder.insert_picture(processed_path)
        except (FileNotFoundError, UnidentifiedImageError) as e:
            logger.error(f"Image processing failed: {str(e)}")
            raise ValueError(f"Image processing failed: {str(e)}")
        except Exception as e:
            logger.error(f"Unexpected error during picture insertion: {str(e)}")
            raise ValueError(f"Unexpected error during picture insertion: {str(e)}")

        # Get placeholder dimensions after picture insertion
        available_width = picture.width
        available_height = picture.height

        # Get original image dimensions directly from the picture object
        image_width, image_height = picture.image.size

        # Calculate aspect ratios
        placeholder_aspect_ratio = float(available_width) / float(available_height)
        image_aspect_ratio = float(image_width) / float(image_height)

        # Store initial position
        pos_left, pos_top = picture.left, picture.top

        picture.crop_top = picture.crop_left = picture.crop_bottom = picture.crop_right = 0

        # Adjust picture dimensions based on aspect ratio comparison
        if placeholder_aspect_ratio > image_aspect_ratio:
            # Placeholder is wider than image - adjust width down while maintaining height
            picture.width = int(image_aspect_ratio * available_height)
            picture.height = available_height
        else:
            # Placeholder is taller than image - adjust height down while maintaining width
            picture.height = int(available_width / image_aspect_ratio)
            picture.width = available_width

        # Center the image within the available space
        picture.left = pos_left + int((available_width - picture.width) / 2)
        picture.top = pos_top + int((available_height - picture.height) / 2)

        # Set the caption
        caption = slide.placeholders[2]
        caption.text_frame.text = caption_text
        caption.text_frame.paragraphs[0].font.size = Pt(14)
        caption.text_frame.paragraphs[0].font.color.rgb = RGBColor(*theme_data["text"])
        caption.text_frame.paragraphs[0].font.name = theme_data["font"]

        return slide

    def add_title_with_content_slide(self, presentation_name: str, title: str, content: str,
                                     theme: str = "modern_blue") -> Slide:
        """
        Add a slide with title and content, applying theme and dynamic text sizing.
        """
        try:
            prs = self.presentations[presentation_name]
        except KeyError as e:
            logger.error(f"Presentation '{presentation_name}' not found")
            raise ValueError(f"Presentation '{presentation_name}' not found")

        theme_data = THEMES.get(theme, THEMES["modern_blue"])
        # Add a slide with title and content
        slide_layout = prs.slide_layouts[self.SLIDE_LAYOUT_TITLE_AND_CONTENT]  # Use layout with title and content
        slide = prs.slides.add_slide(slide_layout)

        # Set the title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*theme_data["text"])
        title_shape.text_frame.paragraphs[0].font.name = theme_data["font"]

        # Set the content
        content_shape = slide.placeholders[1]
        self._add_formatted_bullets(content_shape.text_frame, content, theme_data)

        logger.info(f"Added content slide '{title}' to {presentation_name}")
        return slide

    def add_table_slide(self, presentation_name: str, title: str, headers: List[str], rows: List[List[Union[str, int]]],
                        theme: str = "modern_blue") -> Slide:
        """
        Add a slide with a title and table, applying theme styling.
        """

        try:
            prs = self.presentations[presentation_name]
        except KeyError as e:
            logger.error(f"Presentation '{presentation_name}' not found")
            raise ValueError(f"Presentation '{presentation_name}' not found")
        theme_data = THEMES.get(theme, THEMES["modern_blue"])

        slide_layout = prs.slide_layouts[self.SLIDE_LAYOUT_TITLE_ONLY]
        slide = prs.slides.add_slide(slide_layout)

        # Set the title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*theme_data["text"])
        title_shape.text_frame.paragraphs[0].font.name = theme_data["font"]

        # Calculate table dimensions and position
        num_rows = len(rows) + 1  # +1 for header row
        num_cols = len(headers)

        # Position table in the middle of the slide with some margins
        x = Inches(1)  # Left margin
        y = Inches(2)  # Top margin below title

        # Make table width proportional to the number of columns
        width_per_col = Inches(8 / num_cols)  # Divide available width (8 inches) by number of columns
        height_per_row = Inches(0.4)  # Standard height per row

        # Create table
        shape = slide.shapes.add_table(
            num_rows,
            num_cols,
            x,
            y,
            width_per_col * num_cols,
            height_per_row * num_rows
        )
        table = shape.table

        # Add headers
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)
            # Style header row
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = RGBColor(*theme_data["text"])
            paragraph.font.name = theme_data["font"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*theme_data["accent"])



        # Add data rows
        for row_idx, row_data in enumerate(rows, start=1):
            for col_idx, cell_value in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_value)
                # Style data cells
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = RGBColor(*theme_data["text"])
                paragraph.font.name = theme_data["font"]
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5) if row_idx % 2 == 0 else RGBColor(0xE0, 0xE0, 0xE0)
                logger.info(f"Added table slide '{title}' to {presentation_name}")

        return slide

    def add_title_slide(self, presentation_name: str, title: str, theme: str = "modern_blue") -> Slide:
        """
        Add a title slide with theme styling.
        """
        try:
            prs = self.presentations[presentation_name]
        except KeyError as e:
            logger.error(f"Presentation '{presentation_name}' not found")
            raise ValueError(f"Presentation '{presentation_name}' not found")

        theme_data = THEMES.get(theme, THEMES["modern_blue"])

        # Add a slide with title and content
        slide_layout = prs.slide_layouts[self.SLIDE_LAYOUT_TITLE]
        slide = prs.slides.add_slide(slide_layout)

        # Set the title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*theme_data["text"])
        title_shape.text_frame.paragraphs[0].font.name = theme_data["font"]

        logger.info(f"Added title slide '{title}' to {presentation_name}")
        return slide
    def add_animation(self, presentation_name: str, slide_index: int, effect: str = "fade") -> dict:
        """
        Add an animation effect to a specific slide (Windows-only).
        """
        try:
            import win32com.client
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            prs = powerpoint.Presentations.Open(os.path.abspath(f"{presentation_name}.pptx"))
            try:
                slide = prs.Slides(slide_index + 1)
                slide.Shapes(1).AnimationSettings.EntryEffect = 19 if effect == "fade" else 23  # 19 = Fade, 23 = Zoom
                prs.Save()
                powerpoint.Quit()
                logger.info(f"Added {effect} animation to slide {slide_index} in {presentation_name}")
                return {"status": "animation added"}
            except Exception as e:
                powerpoint.Quit()
                logger.error(f"Failed to add animation: {str(e)}")
                raise ValueError(f"Failed to add animation: {str(e)}")
        except Exception as e:
            logger.error(f"Animation not supported on this platform: {str(e)}")
            return {"status": "animation not supported on this platform"}


